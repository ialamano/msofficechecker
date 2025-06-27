import io
import base64
from flask import Flask, request, jsonify
from flask_cors import CORS
import traceback

# --- Python-docx imports and functions (for Word) ---
from docx import Document
from docx.shared import RGBColor, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement


def add_paragraph_background(paragraph, color_hex):
    """
    Adds a background shade to a paragraph in a Word document.
    color_hex should be an RGB hex string (e.g., "FFFFCC" for light yellow).
    """
    p_pr = paragraph._element.get_or_add_pPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')  # Defines shading type (clear means solid fill)
    shd.set(qn('w:fill'), color_hex)  # Sets the fill color
    p_pr.append(shd)


def check_word_google_docs_compatibility(input_file_base64, original_filename="document.docx"):
    """
    Checks a Word document for compatibility issues when converting to Google Docs
    and generates a report.
    """
    try:
        decoded_bytes = base64.b64decode(input_file_base64)
        doc_bytes_io = io.BytesIO(decoded_bytes)
        doc = Document(doc_bytes_io)

        new_doc = Document()
        new_doc.add_heading('Google Docs Compatibility Report', level=1)
        new_doc.add_paragraph(f'Original File: {original_filename}')
        new_doc.add_paragraph(
            'This report highlights potential compatibility issues when converting this document to Google Docs. '
            'Features like complex formatting, certain fonts, embedded objects (e.g., charts, SmartArt), '
            'and specific layout elements might render differently or not at all.')
        new_doc.add_paragraph('')

        issues_found = []
        original_comments = {c.id: c.text for c in doc.comments}

        # Check for VBA Macros
        if doc.has_macros:
            issues_found.append(
                'VBA Macros: This document contains VBA macros, which are not supported in Google Docs and will be lost upon conversion. Consider converting macro functionality to Google Apps Script if needed.')

        # Check for Comments
        if original_comments:
            issues_found.append(
                'Comments: This document contains comments. While Google Docs supports comments, their appearance and exact positioning might differ after conversion. Original comments will be recreated as new comments in this report and their associated paragraphs highlighted.')

        # Check for Tracked Changes/Revisions
        has_revisions = False
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                run_xml = run._element.xml
                if '<w:ins' in run_xml or '<w:del' in run_xml or '<w:rPrChange' in run_xml:
                    has_revisions = True
                    break
            if has_revisions:
                break
        if has_revisions:
            issues_found.append(
                'Tracked Changes/Revisions: This document contains tracked changes (insertions, deletions, formatting changes). While Google Docs has similar functionality, the way these revisions are displayed or handled (e.g., accepting/rejecting) might differ post-conversion. It is recommended to accept or reject all changes before conversion for a cleaner document.')

        # Add summary of issues
        if issues_found:
            new_doc.add_heading('Summary of Potential Issues', level=2)
            for issue in issues_found:
                new_doc.add_paragraph(f'• {issue}', style='List Bullet')
        else:
            new_doc.add_paragraph(
                'No major compatibility issues (like macros, comments, or tracked changes) were automatically detected. '
                'However, always review the converted document in Google Docs for layout and formatting fidelity.')

        new_doc.add_paragraph('')
        new_doc.add_heading('Original Document Content with Highlights', level=2)
        new_doc.add_paragraph(
            'Sections highlighted in light yellow indicate parts that contained original comments. New comments have been added for clarity regarding these.')
        new_doc.add_paragraph(
            'Note: This tool cannot perfectly simulate Google Docs rendering. Always perform a manual review after conversion.')
        new_doc.add_paragraph('')

        # Copy content, preserving formatting and highlighting comments
        for block in doc.iter_inner_content():
            if isinstance(block, type(doc.paragraphs[0])):  # Check if it's a paragraph
                original_paragraph = block
                new_paragraph = new_doc.add_paragraph()

                # Copy paragraph style
                if original_paragraph.style:
                    try:
                        new_paragraph.style = original_paragraph.style
                    except KeyError:
                        new_paragraph.style = new_doc.styles['Normal']  # Fallback to Normal if style not found

                # Copy runs (text content with formatting)
                for run in original_paragraph.runs:
                    new_run = new_paragraph.add_run(run.text)
                    new_run.bold = run.bold
                    new_run.italic = run.italic
                    new_run.underline = run.underline

                    # Copy font color and size if present
                    if run.font.color.rgb:
                        new_run.font.color.rgb = run.font.color.rgb
                    if run.font.size:
                        new_run.font.size = run.font.size
                    if run.font.name:
                        new_run.font.name = run.font.name

                    # Check for comments and highlight relevant paragraphs
                    run_xml_str = run._element.xml
                    for comment_id, comment_text in original_comments.items():
                        if f'w:commentReference w:id="{comment_id}"' in run_xml_str:
                            add_paragraph_background(new_paragraph, "FFFFCC")  # Light yellow background
                            # Add new comment only once per paragraph
                            if not hasattr(new_paragraph, '_comment_added_by_checker'):
                                new_comment_content = f"Original Comment (ID: {comment_id}): {comment_text}"
                                new_comment = new_doc.add_comment(text=new_comment_content,
                                                                  author='Compatibility Checker', initials='CC')
                                new_run_for_comment_ref = new_paragraph.add_run()
                                new_run_for_comment_ref._element.append(new_comment._element)
                                new_paragraph._comment_added_by_checker = True  # Mark paragraph as having comment added
                            break

            elif isinstance(block, type(doc.tables[0])):  # Check if it's a table
                original_table = block
                new_doc.add_paragraph('')  # Add a space before the table

                # Create a new table and copy content
                new_table = new_doc.add_table(rows=len(original_table.rows), cols=len(original_table.columns))
                for r_idx, row in enumerate(original_table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        new_cell = new_table.cell(r_idx, c_idx)
                        # Clear default paragraph in new cell
                        for p in list(new_cell.paragraphs):
                            new_cell._element.remove(p._element)

                        # Copy paragraphs and runs within each cell
                        for para in cell.paragraphs:
                            new_cell_para = new_cell.add_paragraph()
                            if para.style:
                                try:
                                    new_cell_para.style = para.style
                                except KeyError:
                                    new_cell_para.style = new_doc.styles['Normal']

                            for run in para.runs:
                                new_run = new_cell_para.add_run(run.text)
                                new_run.bold = run.bold
                                new_run.italic = run.italic
                                new_run.underline = run.underline
                                if run.font.color.rgb:
                                    new_run.font.color.rgb = run.font.color.rgb
                                if run.font.size:
                                    new_run.font.size = run.font.size
                                if run.font.name:
                                    new_run.font.name = run.font.name

                new_doc.add_paragraph(
                    'Note: Tables, especially with complex layouts or merged cells, can sometimes have display or formatting issues when converted to Google Docs. Review this table carefully in Google Docs.',
                    style='Intense Quote')

        # Save the new document to bytes and encode for response
        output_bytes_io = io.BytesIO()
        new_doc.save(output_bytes_io)
        output_bytes_encoded = base64.b64encode(output_bytes_io.getvalue()).decode('utf-8')

        return True, output_bytes_encoded, issues_found

    except Exception as e:
        traceback.print_exc()
        return False, "", [
            f"An unexpected error occurred during Word processing: {e}. Please ensure it's a valid .docx file."]


# --- python-pptx imports and functions (for PowerPoint) ---
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def add_warning_textbox(slide, message):
    """Adds a red-bordered textbox with a warning message to a slide."""
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.75)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = message
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)  # Red color for warning

    line = shape.line
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red border
    line.width = Pt(2)


def check_powerpoint_google_slides_compatibility(input_file_base64, original_filename="presentation.pptx"):
    """
    Checks a PowerPoint presentation for compatibility issues when converting to Google Slides
    and generates a report.
    """
    try:
        decoded_bytes = base64.b64decode(input_file_base64)
        pptx_bytes_io = io.BytesIO(decoded_bytes)
        prs = Presentation(pptx_bytes_io)

        new_prs = Presentation()
        new_prs.slide_width = prs.slide_width
        new_prs.slide_height = prs.slide_height

        issues_found = []

        # Create summary slide
        summary_slide_layout = new_prs.slide_layouts[0]  # Title slide layout
        summary_slide = new_prs.slides.add_slide(summary_slide_layout)
        title = summary_slide.shapes.title
        title.text = "Google Slides Compatibility Report"
        subtitle = summary_slide.placeholders[1]
        subtitle.text = f"For: {original_filename}"

        summary_body = summary_slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5)).text_frame
        summary_body.word_wrap = True
        summary_body.add_paragraph().text = "This report highlights potential compatibility issues when converting this presentation to Google Slides. Features like VBA macros, complex animations, specific fonts, and embedded objects might render differently or not at all."
        summary_body.add_paragraph().text = ""

        # Check for VBA Macros
        if prs.has_macros:
            issues_found.append(
                'VBA Macros: This presentation contains VBA macros, which are not supported in Google Slides and will be lost upon conversion. Consider converting macro functionality to Google Apps Script if needed.')

        # Iterate through slides and copy content
        for slide_idx, slide in enumerate(prs.slides):
            try:
                # Try to match layout, fallback to blank layout
                layout = next((l for l in new_prs.slide_layouts if l.name == slide.slide_layout.name),
                              new_prs.slide_layouts[6])  # 6 is blank layout
                new_slide = new_prs.slides.add_slide(layout)
            except Exception:
                new_slide = new_prs.slides.add_slide(new_prs.slide_layouts[6])  # Fallback if layout fails

            # Copy shapes and identify potential issues
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        new_textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                        new_text_frame = new_textbox.text_frame
                        new_text_frame.word_wrap = shape.text_frame.word_wrap
                        for paragraph in shape.text_frame.paragraphs:
                            new_paragraph = new_text_frame.add_paragraph()
                            new_paragraph.text = paragraph.text
                            # Copy basic font formatting
                            if paragraph.font.bold: new_paragraph.font.bold = True
                            if paragraph.font.italic: new_paragraph.font.italic = True
                            if paragraph.font.underline: new_paragraph.font.underline = True
                            if paragraph.font.size: new_paragraph.font.size = paragraph.font.size
                            if paragraph.font.name: new_paragraph.font.name = paragraph.font.name
                            if paragraph.alignment: new_paragraph.alignment = paragraph.alignment

                    elif shape.shape_type == MSO_SHAPE.PICTURE:
                        issues_found.append(
                            f'Slide {slide_idx + 1}: Contains an embedded picture. Image quality or specific effects might differ.')
                    elif not shape.has_text_frame and shape.shape_type != MSO_SHAPE.PICTURE:
                        issues_found.append(
                            f'Slide {slide_idx + 1}: Contains complex graphics or non-standard shapes. Fidelity might be lost upon conversion.')

                except Exception as e:
                    print(f"Error copying shape on slide {slide_idx + 1}: {e}")
                    issues_found.append(
                        f'Slide {slide_idx + 1}: Could not fully copy a shape due to an internal error ({e}). Review this slide carefully.')

            # Check for speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    issues_found.append(
                        f'Slide {slide_idx + 1}: Contains speaker notes. While Google Slides supports notes, their formatting or exact display might differ.')
                    add_warning_textbox(new_slide, f"Original Speaker Notes: {notes_text.strip()[:100]}...")

        # Check for comments in the presentation (overall)
        if prs.comments:  # pptx.Presentation object might not have a direct .comments attribute for collection
            # This check might be more granular per shape or slide if pptx supports it better.
            # Assuming this checks for some top-level comment presence.
            issues_found.append(
                'Comments: This presentation contains comments. While Google Slides supports comments, their appearance and exact positioning might differ after conversion.')

        # Ensure unique issues for final summary display
        issues_found = list(set(issues_found))

        # Add issues summary to the report slide
        if issues_found:
            summary_body.add_paragraph().text = "Summary of Potential Issues:"
            for issue in issues_found:
                summary_body.add_paragraph().text = f'• {issue}'
                summary_body.paragraphs[-1].level = 1  # Indent as a list item
        else:
            summary_body.add_paragraph().text = 'No major compatibility issues (like macros or extensive comments/notes) were automatically detected. However, always review the converted presentation in Google Slides for layout, animation, and formatting fidelity.'

        # Save the new presentation to bytes and encode for response
        output_bytes_io = io.BytesIO()
        new_prs.save(output_bytes_io)
        output_bytes_encoded = base64.b64encode(output_bytes_io.getvalue()).decode('utf-8')

        return True, output_bytes_encoded, issues_found

    except Exception as e:
        traceback.print_exc()
        return False, "", [
            f"An unexpected error occurred during PowerPoint processing: {e}. Please ensure it's a valid .pptx file."]


# --- openpyxl imports and functions (for Excel) ---
from openpyxl import load_workbook, Workbook
from openpyxl.comments import Comment
from openpyxl.styles import PatternFill, Font, Border, Side, Alignment
from openpyxl.utils import get_column_letter, column_index_from_string
from openpyxl.styles.colors import Color  # <--- ADDED: Explicit import for Color object


# Helper function for adding fill to a cell
def add_fill_to_cell(cell, hex_color):
    """Adds a solid background fill to a cell using openpyxl.styles.PatternFill."""
    # openpyxl PatternFill expects a hex string for colors.
    cell.fill = PatternFill(start_color=hex_color, end_color=hex_color, fill_type="solid")


# Helper function for adding a border to a cell (was mentioned in original comment block)
def add_border_to_cell(cell, color="000000", style="thin"):
    """Adds a border to a cell with specified color and style using openpyxl.styles.Border."""
    side = Side(border_style=style, color=color)
    cell.border = Border(left=side, right=side, top=side, bottom=side)


def check_excel_google_sheets_compatibility(input_file_base64, original_filename="document.xlsx"):
    """
    Checks an Excel workbook for compatibility issues when converting to Google Sheets
    and generates a report.
    """
    try:
        decoded_bytes = base64.b64decode(input_file_base64)
        excel_bytes_io = io.BytesIO(decoded_bytes)

        original_wb = load_workbook(excel_bytes_io, keep_vba=True)
        report_wb = Workbook()

        issues_found = []

        summary_ws = report_wb.active
        summary_ws.title = "Compatibility Report"

        # Set up the summary section on the first sheet
        summary_ws['A1'] = "Google Sheets Compatibility Report"
        summary_ws['A1'].font = Font(bold=True, size=16)
        summary_ws['A2'] = f"Original File: {original_filename}"
        summary_ws[
            'A3'] = "This report highlights potential compatibility issues when converting this Excel workbook to Google Sheets."
        summary_ws[
            'A4'] = "Features like VBA macros, complex formulas, certain charts/shapes, and specific formatting might render differently or not at all."

        summary_row = 6  # Starting row for the issues summary list

        # Check for VBA Macros
        if original_wb.vba_archive is not None:
            issues_found.append(
                'VBA Macros: This workbook contains VBA macros, which are not supported in Google Sheets and will be lost upon conversion. Consider converting macro functionality to Google Apps Script if needed.')

        # General warnings for elements openpyxl might not fully parse or are known conversion issues
        issues_found.append(
            'Charts/Complex Graphics: This workbook might contain charts or complex graphic objects. Their appearance, data links, and interactivity might differ significantly or be lost upon conversion.')
        issues_found.append(
            'Embedded Images/Shapes: This workbook might contain embedded images or drawing shapes. Their positioning, scaling, and specific effects might differ upon conversion.')

        # Iterate through each sheet in the original workbook
        for sheet_name in original_wb.sheetnames:
            original_ws = original_wb[sheet_name]
            report_ws = report_wb.create_sheet(title=sheet_name)  # Create a new sheet for the report

            sheet_specific_warnings = []

            # Check for Conditional Formatting
            if original_ws.conditional_formatting:
                sheet_specific_warnings.append(
                    'Conditional Formatting: This sheet uses conditional formatting. While Google Sheets supports some rules, complex or custom rules might not translate perfectly.')
                if not any("Conditional Formatting" in issue for issue in issues_found):
                    issues_found.append(
                        'Conditional Formatting: Some sheets contain conditional formatting rules. Review them carefully after conversion.')

            # Check for Data Validation
            if original_ws.data_validations:
                sheet_specific_warnings.append(
                    'Data Validation: This sheet uses data validation. Simple rules usually transfer, but custom formulas or complex lists might not.')
                if not any("Data Validation" in issue for issue in issues_found):
                    issues_found.append(
                        'Data Validation: Some sheets contain data validation rules. Verify their functionality after conversion.')

            # Add sheet-specific warnings to the report sheet
            if sheet_specific_warnings:
                # Insert rows at the top for warnings
                for _ in range(len(sheet_specific_warnings) + 1):
                    report_ws.insert_rows(1)

                warning_cell_header = report_ws['A1']
                warning_cell_header.value = "Compatibility Warnings for this Sheet:"
                warning_cell_header.font = Font(bold=True,
                                                color=Color(rgb="FF0000"))  # Use openpyxl.styles.colors.Color
                add_fill_to_cell(warning_cell_header, "FFCCCC")  # Light red background
                report_ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=10)  # Merge header

                for idx, warning_msg in enumerate(sheet_specific_warnings):
                    msg_cell = report_ws[f'A{idx + 2}']
                    msg_cell.value = f"- {warning_msg}"
                    msg_cell.font = Font(color=Color(rgb="C80000"))  # Darker red text
                    report_ws.merge_cells(start_row=idx + 2, start_column=1, end_row=idx + 2, end_column=10)

            row_offset = len(sheet_specific_warnings) + 1  # Offset for actual data copying

            # Copy cell data and formatting
            for row_idx, row in enumerate(original_ws.iter_rows()):
                target_row_idx = row_idx + 1 + row_offset

                # Copy row height
                if row_idx + 1 in original_ws.row_dimensions:
                    report_ws.row_dimensions[target_row_idx].height = original_ws.row_dimensions[row_idx + 1].height

                for col_idx, cell in enumerate(row):
                    target_col_letter = get_column_letter(col_idx + 1)
                    report_cell = report_ws[f'{target_col_letter}{target_row_idx}']

                    report_cell.value = cell.value  # Copy cell value

                    # Copy cell styles (font, fill, alignment, border)
                    if cell.has_style:
                        if cell.font:
                            font_color_value = None
                            if cell.font.color:
                                try:
                                    # Safely extract RGB hex string for openpyxl Font color
                                    if isinstance(cell.font.color, Color) and cell.font.color.rgb and len(
                                            cell.font.color.rgb) in [6, 8]:
                                        font_color_value = cell.font.color.rgb
                                    # Fallback for unexpected color types (e.g., from pptx if mixed up, though unlikely here)
                                    elif hasattr(cell.font.color, 'rgb') and hasattr(cell.font.color.rgb, 'r'):
                                        rgb = cell.font.color.rgb
                                        font_color_value = f"{rgb.r:02X}{rgb.g:02X}{rgb.b:02X}"
                                except Exception:
                                    font_color_value = '000000'  # Default to black if conversion fails

                            report_cell.font = Font(name=cell.font.name, size=cell.font.size,
                                                    bold=cell.font.bold, italic=cell.font.italic,
                                                    underline=cell.font.underline, strike=cell.font.strike,
                                                    color=font_color_value)  # Use the extracted hex string or None
                        if cell.fill:
                            # openpyxl fill colors also expect hex strings.
                            report_cell.fill = PatternFill(start_color=cell.fill.start_color,
                                                           end_color=cell.fill.end_color, fill_type=cell.fill.fill_type)
                        if cell.alignment:
                            report_cell.alignment = Alignment(horizontal=cell.alignment.horizontal,
                                                              vertical=cell.alignment.vertical,
                                                              wrap_text=cell.alignment.wrap_text,
                                                              shrink_to_fit=cell.alignment.shrink_to_fit,
                                                              indent=cell.alignment.indent,
                                                              text_rotation=cell.alignment.text_rotation,
                                                              readingOrder=cell.alignment.readingOrder)
                        if cell.border:
                            # openpyxl border colors also expect hex strings.
                            report_cell.border = Border(left=cell.border.left, right=cell.border.right,
                                                        top=cell.border.top, bottom=cell.border.bottom,
                                                        diagonalUp=cell.border.diagonalUp,
                                                        diagonalDown=cell.border.diagonalDown,
                                                        outline=cell.border.outline)

                    # Copy column width
                    if col_idx + 1 in original_ws.column_dimensions:
                        report_ws.column_dimensions[target_col_letter].width = original_ws.column_dimensions[
                            get_column_letter(col_idx + 1)].width

                    # Check for cell comments
                    if cell.comment:
                        add_fill_to_cell(report_cell, "FFFFCC")  # Highlight cells with comments
                        report_cell.comment = Comment(f"Original Comment: {cell.comment.text}", "Compatibility Checker")
                        if not any("Comments" in issue for issue in issues_found):
                            issues_found.append(
                                'Comments: Some cells contain comments. While Google Sheets supports comments, their appearance and exact positioning might differ.')

                    # Check for formulas
                    if cell.data_type == 'f':  # 'f' denotes a formula cell
                        add_fill_to_cell(report_cell, "FFFFCC")  # Highlight formulas
                        if not any("Formulas" in issue for issue in issues_found):
                            issues_found.append(
                                'Formulas: This workbook contains formulas. Most common functions transfer, but complex array formulas or Excel-specific functions might break. Review formulas after conversion.')

            # Handle merged cells
            for merged_range_obj in original_ws.merged_cells.ranges:
                start_col = merged_range_obj.min_col
                start_row = merged_range_obj.min_row
                end_col = merged_range_obj.max_col
                end_row = merged_range_obj.max_row

                # Adjust merged cell range for row offset
                new_start_row = start_row + row_offset
                new_end_row = end_row + row_offset

                new_merge_range = f"{get_column_letter(start_col)}{new_start_row}:{get_column_letter(end_col)}{new_end_row}"
                report_ws.merge_cells(new_merge_range)

                # Highlight the top-left cell of the merged range
                top_left_cell_in_report = report_ws[f"{get_column_letter(start_col)}{new_start_row}"]
                add_fill_to_cell(top_left_cell_in_report, "FFFFCC")  # Highlight merged cells

                if not any("Merged Cells" in issue for issue in issues_found):
                    issues_found.append(
                        'Merged Cells: This workbook contains merged cells. While Google Sheets supports merging, visual layout might differ subtly.')

        # Final summary on the first sheet
        summary_ws['A' + str(summary_row)].value = "Summary of Potential Issues:"
        summary_ws['A' + str(summary_row)].font = Font(bold=True)
        summary_row += 1

        if issues_found:
            issues_found = list(set(issues_found))  # Ensure unique issues
            for issue in issues_found:
                summary_ws['A' + str(summary_row)].value = f"• {issue}"
                summary_ws['A' + str(summary_row)].font = Font(size=10)
                summary_row += 1
        else:
            summary_ws['A' + str(
                summary_row)].value = "No major compatibility issues were automatically detected. However, always review the converted spreadsheet in Google Sheets for layout and formatting fidelity."
            summary_ws['A' + str(summary_row)].font = Font(italic=True,
                                                           color=Color(rgb="646464"))  # Grey italic for no issues

        # Save the new workbook to bytes and encode for response
        output_bytes_io = io.BytesIO()
        report_wb.save(output_bytes_io)
        output_bytes_encoded = base64.b64encode(output_bytes_io.getvalue()).decode('utf-8')

        return True, output_bytes_encoded, issues_found

    except Exception as e:
        traceback.print_exc()
        return False, "", [
            f"An unexpected error occurred during Excel processing: {e}. Please ensure it's a valid .xlsx or .xlsm file. Details: {str(e)}"]


# --- Flask App Setup ---
app = Flask(__name__)
CORS(app)  # Allow cross-origin requests for development


# Root route for server health check
@app.route('/', methods=['GET'])
def home():
    """Simple health check endpoint."""
    return "Unified Office File Compatibility Checker Backend is running!"


# Main compatibility check endpoint
@app.route('/check-compatibility', methods=['POST'])
def check_compatibility_route():
    """
    Handles POST requests for file compatibility checks.
    Expects file_base64, filename, and file_type in the JSON payload.
    """
    try:
        data = request.get_json()
        file_base64 = data.get('file_base64')
        filename = data.get('filename', 'document.file')
        file_type = data.get('file_type')  # e.g., 'docx', 'pptx', 'xlsx'

        if not file_base64:
            return jsonify({'success': False, 'error': 'No file data provided.'}), 400

        # Determine which checker function to call based on file_type
        if file_type == 'docx':
            success, output_base64, issues = check_word_google_docs_compatibility(file_base64, filename)
        elif file_type == 'pptx':
            success, output_base64, issues = check_powerpoint_google_slides_compatibility(file_base64, filename)
        elif file_type == 'xlsx' or file_type == 'xlsm':
            success, output_base64, issues = check_excel_google_sheets_compatibility(file_base64, filename)
        else:
            return jsonify({'success': False, 'error': 'Unsupported file type provided.'}), 400

        if success:
            return jsonify({
                'success': True,
                'output_file_base64': output_base64,
                'issues_found': issues
            })
        else:
            # If checker function returns False, issues list contains the error message
            return jsonify({
                'success': False,
                'error': issues[0] if issues else 'Unknown error during processing.'
            }), 500

    except Exception as e:
        # Catch any unexpected errors during the request handling itself
        traceback.print_exc()  # Print full traceback to console for debugging on server
        return jsonify({'success': False, 'error': str(e)}), 500


# Entry point for running the Flask app
if __name__ == '__main__':
    # Run the app in debug mode locally, accessible from any IP on port 5000
    app.run(debug=True, host='0.0.0.0', port=5000)

